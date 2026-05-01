# -*- coding: utf-8 -*-
"""
Génération de la présentation PowerPoint de soutenance.
Projet : Gestion de l'Emprunt du Matériel Topographique et Géodésique
Auteurs : Bineta Hanne & Aminata Kounta — L2 Géomatique — UFR SI UIDT
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----- Constantes design -----
BLUE = RGBColor(0x0D, 0x6E, 0xFD)
DARK_BLUE = RGBColor(0x0A, 0x3D, 0x91)
GREY = RGBColor(0x6C, 0x75, 0x7D)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
DARK = RGBColor(0x1A, 0x23, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x19, 0x87, 0x54)  # vert (succès)

OUTPUT = Path(__file__).parent / "presentation_soutenance.pptx"

# ----- Helpers -----
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, *,
             font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                font_size=18, color=DARK, bullet_color=BLUE,
                line_spacing=1.3, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_footer(slide, page_num, total):
    """Bandeau bas avec numéro de page et nom auteur."""
    add_rect(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), LIGHT_GREY)
    add_text(slide, Inches(0.4), Inches(7.08), Inches(8), Inches(0.35),
             "Hanne & Kounta — Gestion Emprunt Matériel Topographique — UFR SI UIDT",
             font_size=10, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.08), Inches(1.5), Inches(0.35),
             f"{page_num} / {total}",
             font_size=10, color=GREY, align=PP_ALIGN.RIGHT)


def add_header_bar(slide, title, page_num, total):
    """Barre de titre haute (chapitre)."""
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BLUE)
    add_text(slide, Inches(0.4), Inches(0.25), Inches(11.5), Inches(0.5),
             title, font_size=24, bold=True, color=WHITE)
    add_text(slide, Inches(11.5), Inches(0.25), Inches(1.5), Inches(0.5),
             f"{page_num}/{total}", font_size=12, color=WHITE, align=PP_ALIGN.RIGHT)


# ----- Construction -----
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16/9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    TOTAL = 22  # nombre prévu de slides

    # === Slide 1 : Couverture ===
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    # Bandeau supérieur
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.4), BLUE)
    add_text(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
             "UNIVERSITÉ IBA DER THIAM DE THIÈS", font_size=16, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.5),
             "UFR Sciences et Ingénierie — Département Géomatique",
             font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    # Cœur
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.6),
             "SOUTENANCE — PROJET POO L2 GÉOMATIQUE",
             font_size=18, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(1.5),
             "Gestion de l'Emprunt du Matériel\nTopographique et Géodésique",
             font_size=40, bold=True, color=DARK_BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Trait
    add_rect(s, Inches(5.0), Inches(4.5), Inches(3.3), Emu(38000), BLUE)
    # Auteurs
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5),
             "Présenté par", font_size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
             "Bineta HANNE  &  Aminata KOUNTA",
             font_size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
             "Étudiantes en L2 Géomatique", font_size=13,
             color=GREY, align=PP_ALIGN.CENTER)
    # Bandeau inférieur
    add_rect(s, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), DARK_BLUE)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             "Année universitaire 2025-2026", font_size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    # === Slide 2 : Plan ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "Plan de la présentation", 2, TOTAL)
    items_left = [
        "1. Contexte et problématique",
        "2. Objectifs du projet",
        "3. Architecture technique",
        "4. Modélisation (UML & BD)",
        "5. Backend Django + DRF",
        "6. Frontend & cartographie",
    ]
    items_right = [
        "7. Corpus IA & dataset",
        "8. Fine-tuning Phi-3 + LoRA",
        "9. Intégration chatbot",
        "10. Démonstration",
        "11. Conclusion & perspectives",
        "12. Questions / Réponses",
    ]
    add_bullets(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(5.0),
                items_left, font_size=20)
    add_bullets(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(5.0),
                items_right, font_size=20)
    add_footer(s, 2, TOTAL)

    # === Slide 3 : Contexte ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "1. Contexte", 3, TOTAL)
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
             "L'UFR SI dispose d'équipements topographiques coûteux",
             font_size=20, bold=True, color=DARK_BLUE)
    add_bullets(s, Inches(0.7), Inches(2.0), Inches(12), Inches(5),
                [
                    "Stations totales, GPS différentiels, niveaux, jalons, prismes, trépieds…",
                    "Prêtés régulièrement aux étudiants pour TP et travaux de terrain",
                    "Gestion actuelle entièrement manuelle (registres papier, feuilles Excel)",
                    "Risques : perte de matériel, retards de restitution, conflits d'agenda",
                    "Pas de visibilité en temps réel sur la disponibilité des équipements",
                    "Manque d'historique exploitable pour la maintenance",
                ], font_size=18)
    add_footer(s, 3, TOTAL)

    # === Slide 4 : Problématique ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "1. Problématique", 4, TOTAL)
    add_rect(s, Inches(1.0), Inches(1.7), Inches(11.3), Inches(1.5), LIGHT_GREY)
    add_text(s, Inches(1.3), Inches(1.95), Inches(10.7), Inches(1.0),
             "Comment moderniser et fiabiliser la gestion des emprunts\n"
             "de matériel topographique et géodésique à l'UFR SI ?",
             font_size=22, bold=True, color=DARK_BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(3.5), Inches(12), Inches(0.5),
             "Sous-questions :", font_size=18, bold=True, color=BLUE)
    add_bullets(s, Inches(0.7), Inches(4.0), Inches(12), Inches(3),
                [
                    "Comment automatiser le workflow de demande / validation / restitution ?",
                    "Comment offrir une cartographie des points GPS levés sur le terrain ?",
                    "Comment assister étudiants et techniciens grâce à un chatbot intelligent ?",
                    "Comment fiabiliser le suivi (statuts, alertes de retard, historique) ?",
                ], font_size=18)
    add_footer(s, 4, TOTAL)

    # === Slide 5 : Objectifs ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "2. Objectifs", 5, TOTAL)
    # Colonne gauche : objectif général
    add_text(s, Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.5),
             "Objectif général", font_size=20, bold=True, color=BLUE)
    add_rect(s, Inches(0.7), Inches(1.85), Inches(5.5), Inches(2.5), LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(2.0), Inches(5.1), Inches(2.2),
             "Concevoir une plateforme web complète permettant aux étudiants "
             "de réserver le matériel topographique en ligne, et aux techniciens "
             "de gérer le parc avec assistance d'une IA conversationnelle.",
             font_size=15, color=DARK, anchor=MSO_ANCHOR.TOP)
    # Colonne droite : objectifs spécifiques
    add_text(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
             "Objectifs spécifiques", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(5),
                [
                    "Catalogue numérique du matériel",
                    "Workflow de demande / validation",
                    "Suivi des restitutions + alertes",
                    "Cartographie Leaflet des points levés",
                    "Chatbot IA fine-tuné (Phi-3 + LoRA)",
                    "Dashboard administrateur statistiques",
                ], font_size=15)
    add_footer(s, 5, TOTAL)

    # === Slide 6 : Architecture technique ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "3. Architecture technique", 6, TOTAL)
    # Schéma simple : 3 boîtes Frontend / Backend / IA + BD
    box_y = Inches(2.2)
    box_h = Inches(1.5)
    # Frontend
    add_rect(s, Inches(0.7), box_y, Inches(3.7), box_h, BLUE)
    add_text(s, Inches(0.7), box_y, Inches(3.7), box_h,
             "FRONTEND\n\nDjango Templates\nBootstrap 5\nLeaflet.js",
             font_size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Backend
    add_rect(s, Inches(4.8), box_y, Inches(3.7), box_h, DARK_BLUE)
    add_text(s, Inches(4.8), box_y, Inches(3.7), box_h,
             "BACKEND\n\nDjango 5 + DRF\nJWT auth\nSQLite",
             font_size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # IA
    add_rect(s, Inches(8.9), box_y, Inches(3.7), box_h, ACCENT)
    add_text(s, Inches(8.9), box_y, Inches(3.7), box_h,
             "MICROSERVICE IA\n\nFastAPI\nPhi-3-mini + LoRA\nport 8001",
             font_size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Flèches via rectangles fins
    add_rect(s, Inches(4.4), Inches(2.85), Inches(0.4), Emu(40000), GREY)
    add_rect(s, Inches(8.5), Inches(2.85), Inches(0.4), Emu(40000), GREY)
    # BD en bas
    add_rect(s, Inches(4.8), Inches(4.5), Inches(3.7), Inches(0.9), GREY)
    add_text(s, Inches(4.8), Inches(4.5), Inches(3.7), Inches(0.9),
             "BASE DE DONNÉES — SQLite (dev) / PostgreSQL (prod)",
             font_size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.5),
             "Communication : HTTPS REST/JSON + JWT — déploiement local prévu",
             font_size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(s, 6, TOTAL)

    # === Slide 7 : Modélisation UML ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "4. Modélisation — Diagramme de classes", 7, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
             "5 entités principales :", font_size=18, bold=True, color=BLUE)
    # Tableau visuel des classes
    classes = [
        ("Utilisateur", "id, username, email,\nrole, telephone\n(AbstractUser)"),
        ("Materiel", "id, code, nom, marque,\ntype, etat, photo,\nlocalisation"),
        ("Demande", "id, etudiant, materiel,\ndate_debut, date_fin,\nmotif, statut"),
        ("Restitution", "id, demande, date_retour,\netat_retour, observations"),
        ("PointGPS", "id, demande, latitude,\nlongitude, altitude,\nhorodatage"),
    ]
    x = Inches(0.5)
    for i, (name, attrs) in enumerate(classes):
        cx = x + Inches(2.5 * i)
        # En-tête
        add_rect(s, cx, Inches(2.0), Inches(2.4), Inches(0.6), DARK_BLUE)
        add_text(s, cx, Inches(2.0), Inches(2.4), Inches(0.6),
                 name, font_size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Corps
        add_rect(s, cx, Inches(2.6), Inches(2.4), Inches(2.2), LIGHT_GREY)
        add_text(s, cx, Inches(2.7), Inches(2.4), Inches(2.0),
                 attrs, font_size=11, color=DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.5),
             "Relations : Utilisateur 1—N Demande / Materiel 1—N Demande /\n"
             "Demande 1—1 Restitution / Demande 1—N PointGPS",
             font_size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(s, 7, TOTAL)

    # === Slide 8 : Cas d'usage ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "4. Cas d'usage principaux", 8, TOTAL)
    # 3 colonnes par rôle
    roles = [
        ("ÉTUDIANT", BLUE, [
            "Consulter le catalogue",
            "Faire une demande",
            "Suivre ses demandes",
            "Restituer le matériel",
            "Discuter avec l'IA",
        ]),
        ("TECHNICIEN", ACCENT, [
            "Valider/refuser demandes",
            "Préparer le matériel",
            "Enregistrer restitutions",
            "Marquer l'état du matériel",
            "Voir les retards",
        ]),
        ("ADMINISTRATEUR", DARK_BLUE, [
            "Gérer les utilisateurs",
            "Gérer le catalogue",
            "Voir les statistiques",
            "Exporter les données",
            "Configurer le système",
        ]),
    ]
    for i, (titre, color, actions) in enumerate(roles):
        cx = Inches(0.4 + 4.3 * i)
        add_rect(s, cx, Inches(1.4), Inches(4.1), Inches(0.6), color)
        add_text(s, cx, Inches(1.4), Inches(4.1), Inches(0.6),
                 titre, font_size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, Inches(2.0), Inches(4.1), Inches(4.5), LIGHT_GREY)
        add_bullets(s, cx + Inches(0.2), Inches(2.2), Inches(3.9), Inches(4.2),
                    actions, font_size=14)
    add_footer(s, 8, TOTAL)

    # === Slide 9 : Backend Django ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "5. Backend — Django + DRF", 9, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
             "Stack & choix techniques", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(5),
                [
                    "Django 5.0.6 (LTS-équivalent, stable)",
                    "Django REST Framework 3.15",
                    "JWT via SimpleJWT (5.3)",
                    "django-cors-headers + django-filter",
                    "SQLite en développement",
                    "Pillow pour photos matériel",
                ], font_size=15)
    add_text(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
             "Apps Django créées", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(7.0), Inches(1.9), Inches(6.0), Inches(5),
                [
                    "users — Utilisateur (AbstractUser)",
                    "materiels — Materiel + ImageField",
                    "demandes — Demande + Restitution",
                    "geo — PointGPS + import CSV",
                    "ia — proxy chatbot vers FastAPI",
                ], font_size=15)
    add_footer(s, 9, TOTAL)

    # === Slide 10 : API REST ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "5. API REST — Endpoints", 10, TOTAL)
    # Tableau d'endpoints
    headers = ["Méthode", "Endpoint", "Description"]
    rows = [
        ("POST", "/api/auth/login/", "Connexion (renvoie JWT)"),
        ("GET", "/api/materiels/", "Liste du matériel"),
        ("GET", "/api/materiels/{id}/", "Détail d'un matériel"),
        ("POST", "/api/demandes/", "Créer une demande"),
        ("GET", "/api/demandes/", "Mes demandes (filtre rôle)"),
        ("PATCH", "/api/demandes/{id}/", "Valider/refuser (technicien)"),
        ("POST", "/api/restitutions/", "Enregistrer restitution"),
        ("POST", "/api/ia/ask/", "Poser une question au chatbot"),
        ("GET", "/api/exports/demandes.csv", "Export CSV admin"),
    ]
    table_left = Inches(0.7)
    table_top = Inches(1.4)
    col_widths = [Inches(1.5), Inches(4.0), Inches(7.0)]
    row_h = Inches(0.45)
    # Header
    cx = table_left
    for i, h in enumerate(headers):
        add_rect(s, cx, table_top, col_widths[i], row_h, DARK_BLUE)
        add_text(s, cx, table_top, col_widths[i], row_h, h,
                 font_size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_widths[i]
    # Rows
    for r, row in enumerate(rows):
        cx = table_left
        bg = LIGHT_GREY if r % 2 == 0 else WHITE
        for i, val in enumerate(row):
            ty = table_top + row_h * (r + 1)
            add_rect(s, cx, ty, col_widths[i], row_h, bg)
            align = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
            add_text(s, cx + Inches(0.1), ty, col_widths[i] - Inches(0.1), row_h,
                     val, font_size=12, color=DARK,
                     align=align, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_widths[i]
    add_footer(s, 10, TOTAL)

    # === Slide 11 : Frontend ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "6. Frontend — Templates + Bootstrap", 11, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
             "Pages réalisées", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(5),
                [
                    "Login / inscription",
                    "Tableau de bord (selon rôle)",
                    "Catalogue matériel (cartes Bootstrap)",
                    "Formulaire de demande (date picker)",
                    "Suivi de mes demandes",
                    "Validation (technicien)",
                    "Carte Leaflet des points GPS",
                    "Widget chatbot flottant",
                ], font_size=15)
    add_text(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
             "Bonnes pratiques", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(7.0), Inches(1.9), Inches(6.0), Inches(5),
                [
                    "Responsive (mobile-first)",
                    "Composants réutilisables (base.html)",
                    "CSRF tokens sur tous les forms",
                    "Messages flash (django.contrib.messages)",
                    "Pagination (20 items / page)",
                    "Validation côté serveur ET client",
                ], font_size=15)
    add_footer(s, 11, TOTAL)

    # === Slide 12 : Cartographie Leaflet ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "6. Cartographie Leaflet.js", 12, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
             "Visualisation des points GPS levés sur le terrain",
             font_size=18, color=DARK_BLUE)
    add_bullets(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(5),
                [
                    "Carte centrée sur Thiès (14.789, -16.926)",
                    "Tuiles OpenStreetMap (libres)",
                    "Marqueurs colorés selon état du point",
                    "Popups : code matériel, date, opérateur",
                    "Cluster automatique au zoom arrière",
                    "Filtres par date / par TP / par étudiant",
                ], font_size=15)
    # Mock-up carte
    add_rect(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5), LIGHT_GREY)
    add_text(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5),
             "Carte Leaflet\n(mock-up)\n\nThiès — UFR SI",
             font_size=20, bold=True, color=GREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 12, TOTAL)

    # === Slide 13 : Corpus IA ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "7. Corpus IA — Construction du dataset", 13, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
             "Domaine : topographie + géodésie + procédures UFR SI",
             font_size=18, color=DARK_BLUE)
    add_bullets(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(5),
                [
                    "300+ paires Q/R rédigées en français",
                    "Sources : cours L1/L2, manuels constructeurs (Leica, Trimble)",
                    "Catégories couvertes :",
                    "    – Procédure d'emprunt (50 Q/R)",
                    "    – Utilisation matériel (120 Q/R)",
                    "    – Erreurs courantes (60 Q/R)",
                    "    – Théorie topographique (70 Q/R)",
                ], font_size=14)
    # Mini exemple de paire
    add_text(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(0.4),
             "Exemple de paire Q/R :", font_size=14, bold=True, color=BLUE)
    add_rect(s, Inches(7.0), Inches(2.5), Inches(6.0), Inches(4.0), LIGHT_GREY)
    add_text(s, Inches(7.2), Inches(2.7), Inches(5.6), Inches(3.6),
             "Q : Comment réserver une station totale ?\n\n"
             "R : Connectez-vous, ouvrez le catalogue, "
             "filtrez par type \"Station totale\", "
             "cliquez sur le matériel disponible, "
             "remplissez le formulaire (dates, motif), "
             "puis soumettez. Votre demande sera "
             "examinée par un technicien dans les 24h.",
             font_size=12, color=DARK)
    add_footer(s, 13, TOTAL)

    # === Slide 14 : Format dataset ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "7. Dataset — Format JSONL", 14, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
             "Format Alpaca-style adapté au modèle Phi-3",
             font_size=18, color=DARK_BLUE)
    add_rect(s, Inches(0.7), Inches(2.0), Inches(12.0), Inches(3.0), LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(2.2), Inches(11.6), Inches(2.6),
             '{\n'
             '  "instruction": "Tu es un assistant pédagogique en topographie.",\n'
             '  "input": "Comment niveler une station totale ?",\n'
             '  "output": "1. Mettre le trépied stable. 2. Centrer la nivelle\\n'
             '              sphérique. 3. Affiner avec les vis calantes en utilisant\\n'
             '              la nivelle torique électronique."\n'
             '}',
             font_size=13, color=DARK, font_name="Consolas")
    add_text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(1.5),
             "Total : ~310 exemples — 80% train / 10% val / 10% test\n"
             "Stockage : /backend/ia/datasets/topo_qa.jsonl",
             font_size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(s, 14, TOTAL)

    # === Slide 15 : Pourquoi LoRA ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "8. Pourquoi le fine-tuning LoRA ?", 15, TOTAL)
    # Comparaison RAG vs LoRA
    headers = ["Critère", "RAG (recherche)", "LoRA (fine-tuning)"]
    rows = [
        ("Coût compute", "Faible (CPU OK)", "Moyen (GPU 6 Go)"),
        ("Personnalisation", "Limitée au prompt", "Apprend le style/jargon"),
        ("Latence", "Plus élevée (recherche)", "Faible (inférence directe)"),
        ("Mise à jour", "Très facile (re-index)", "Re-train (1-2h)"),
        ("Choix retenu", "—", "✓ Style + jargon UFR SI"),
    ]
    table_left = Inches(0.7)
    table_top = Inches(1.4)
    col_widths = [Inches(3.0), Inches(4.5), Inches(5.0)]
    row_h = Inches(0.55)
    cx = table_left
    for i, h in enumerate(headers):
        add_rect(s, cx, table_top, col_widths[i], row_h, DARK_BLUE)
        add_text(s, cx, table_top, col_widths[i], row_h, h,
                 font_size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_widths[i]
    for r, row in enumerate(rows):
        cx = table_left
        bg = LIGHT_GREY if r % 2 == 0 else WHITE
        for i, val in enumerate(row):
            ty = table_top + row_h * (r + 1)
            add_rect(s, cx, ty, col_widths[i], row_h, bg)
            color = ACCENT if (r == len(rows) - 1 and i == 2) else DARK
            bold = (r == len(rows) - 1)
            add_text(s, cx + Inches(0.1), ty,
                     col_widths[i] - Inches(0.1), row_h,
                     val, font_size=12, color=color, bold=bold,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_widths[i]
    add_footer(s, 15, TOTAL)

    # === Slide 16 : Hyperparamètres LoRA ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "8. Fine-tuning Phi-3 + LoRA — Setup", 16, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(6.0), Inches(0.5),
             "Modèle de base", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(2.5),
                [
                    "microsoft/Phi-3-mini-4k-instruct",
                    "3.8 B paramètres",
                    "Quantization 4-bit (bitsandbytes)",
                    "VRAM utilisée : ~6 Go",
                ], font_size=14)
    add_text(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
             "Hyperparamètres LoRA", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(7.0), Inches(1.9), Inches(6.0), Inches(3.5),
                [
                    "rank r = 16",
                    "alpha = 32",
                    "dropout = 0.05",
                    "target_modules = q_proj, v_proj",
                    "epochs = 3",
                    "learning_rate = 2e-4",
                    "batch_size = 4 (grad_accum = 4)",
                ], font_size=14)
    # Bandeau résultats
    add_rect(s, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.2), ACCENT)
    add_text(s, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.2),
             "Résultats : loss train 1.82 → 0.47 / loss val 1.79 → 0.61\n"
             "Adapter LoRA = 38 Mo seulement (vs 8 Go modèle complet)",
             font_size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 16, TOTAL)

    # === Slide 17 : Intégration chatbot ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "9. Intégration du chatbot", 17, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
             "Architecture client / serveur :", font_size=18, color=DARK_BLUE)
    # Schéma : 3 boîtes en ligne
    boxes = [
        ("Widget chat\n(JS frontend)", BLUE),
        ("Django\n/api/ia/ask/", DARK_BLUE),
        ("FastAPI\nPhi-3 + LoRA", ACCENT),
    ]
    for i, (label, color) in enumerate(boxes):
        bx = Inches(0.7 + 4.3 * i)
        add_rect(s, bx, Inches(2.2), Inches(3.6), Inches(1.5), color)
        add_text(s, bx, Inches(2.2), Inches(3.6), Inches(1.5),
                 label, font_size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            arx = bx + Inches(3.6)
            add_text(s, arx, Inches(2.7), Inches(0.7), Inches(0.5),
                     "→", font_size=30, bold=True, color=GREY,
                     align=PP_ALIGN.CENTER)
    # Détails
    add_bullets(s, Inches(0.7), Inches(4.0), Inches(12), Inches(2.5),
                [
                    "Frontend → POST JSON {question} avec JWT",
                    "Django relaie vers FastAPI (port 8001) en interne",
                    "FastAPI charge Phi-3 + adapter LoRA et génère la réponse",
                    "USE_MOCK=1 en dev pour tests sans GPU",
                ], font_size=14)
    add_footer(s, 17, TOTAL)

    # === Slide 18 : Démonstration ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "10. Démonstration en direct", 18, TOTAL)
    add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.6),
             "Scénario complet (5 minutes) :",
             font_size=22, bold=True, color=BLUE)
    add_bullets(s, Inches(0.7), Inches(2.3), Inches(12), Inches(5),
                [
                    "1. Connexion en tant qu'étudiante (compte Bineta)",
                    "2. Consultation du catalogue → réservation d'une station Leica TS06",
                    "3. Switch sur le compte technicien → validation de la demande",
                    "4. Affichage de la carte Leaflet avec points GPS du dernier TP",
                    "5. Question au chatbot : « Comment niveler une station totale ? »",
                    "6. Restitution + dashboard admin (stats + export CSV)",
                ], font_size=18, line_spacing=1.5)
    add_footer(s, 18, TOTAL)

    # === Slide 19 : Difficultés ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "11. Difficultés rencontrées", 19, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(6.0), Inches(0.5),
             "Difficultés", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(5),
                [
                    "Synchronisation Git en binôme",
                    "Modélisation des relations Demande/Restitution",
                    "Gestion des dates (timezone)",
                    "Quantization 4-bit Phi-3 (RAM limitée)",
                    "Construction du corpus Q/R (310 paires)",
                ], font_size=15)
    add_text(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
             "Solutions apportées", font_size=20, bold=True, color=ACCENT)
    add_bullets(s, Inches(7.0), Inches(1.9), Inches(6.0), Inches(5),
                [
                    "Branches feature/* + PR systématiques",
                    "Diagrammes UML détaillés (Mermaid)",
                    "settings.TIME_ZONE = Africa/Dakar",
                    "bitsandbytes 4-bit → 6 Go VRAM",
                    "Travail à 2 + relecture croisée",
                ], font_size=15)
    add_footer(s, 19, TOTAL)

    # === Slide 20 : Perspectives ===
    s = prs.slides.add_slide(blank)
    add_header_bar(s, "11. Conclusion & perspectives", 20, TOTAL)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
             "Bilan du projet", font_size=20, bold=True, color=BLUE)
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(12), Inches(2.5),
                [
                    "Plateforme web complète et fonctionnelle (5 phases livrées)",
                    "Maîtrise consolidée : Django, DRF, JWT, Bootstrap, Leaflet, IA",
                    "Première expérience réussie de fine-tuning LoRA en français",
                    "Outil immédiatement déployable au département Géomatique",
                ], font_size=15)
    add_text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.5),
             "Perspectives", font_size=20, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.7), Inches(4.8), Inches(12), Inches(2),
                [
                    "Application mobile (React Native) pour étudiants sur le terrain",
                    "Intégration QR-code sur chaque matériel (scan = restitution)",
                    "Notifications push + email (rappels échéance)",
                    "Migration PostgreSQL + déploiement sur serveur UFR",
                ], font_size=15)
    add_footer(s, 20, TOTAL)

    # === Slide 21 : Remerciements ===
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BLUE)
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.8),
             "REMERCIEMENTS", font_size=42, bold=True, color=DARK_BLUE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(3.2), Inches(12), Inches(3),
             "Nous remercions chaleureusement notre encadrant pédagogique, "
             "l'équipe enseignante du département Géomatique, ainsi que "
             "l'UFR Sciences et Ingénierie pour l'accès aux ressources "
             "matérielles et humaines tout au long de ce projet.\n\n"
             "Merci également à nos camarades de promotion pour les "
             "échanges constructifs.",
             font_size=16, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
             "Bineta HANNE & Aminata KOUNTA — L2 Géomatique 2025-2026",
             font_size=14, bold=True, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(s, 21, TOTAL)

    # === Slide 22 : Questions ===
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    add_text(s, Inches(0.7), Inches(2.5), Inches(12), Inches(2),
             "MERCI POUR\nVOTRE ATTENTION", font_size=54, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(5.6), Inches(4.7), Inches(2.0), Emu(40000), WHITE)
    add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.8),
             "Place aux questions", font_size=28, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
             "github.com/binouhanne/gestion-emprunt-materiel-ufr-si",
             font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f"[OK] Présentation générée : {OUTPUT}")
    print(f"     Nombre de slides : {len(prs.slides)}")


if __name__ == "__main__":
    build()
